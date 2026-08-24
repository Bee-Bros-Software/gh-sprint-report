"""PowerPoint sprint-review deck generation.

Builds a self-contained ``.pptx`` from computed metrics. Charts are created as
native PowerPoint charts rather than rendered images, so numbers stay editable
and the deck can be adjusted live during review.

Slide order:

1. Title
2. Sprint summary (headline figures)
3. Burndown
4. Burnup
4. Velocity and rolling average
5. Commitment vs completed
6. Work mix (planned / unplanned / carryover)
7. Rolling forward (incomplete work)
8. Forecast

Example:
    >>> import tempfile, pathlib
    >>> from sprint_report.models import SprintMetrics
    >>> builder = DeckBuilder(project_title="Acme Platform")
    >>> with tempfile.TemporaryDirectory() as tmp:
    ...     out = builder.build(
    ...         current=SprintMetrics("Sprint 14", committed_points=40,
    ...                               completed_points=34),
    ...         history=[], burndown_points=[], carryover=[],
    ...         output_path=pathlib.Path(tmp) / "deck.pptx",
    ...     )
    ...     out.exists()
    True
"""

from __future__ import annotations

from collections.abc import Sequence
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .metrics import ScopeChange, forecast_sprints, rolling_average
from .models import BurndownPoint, ProjectItem, SprintMetrics

__all__ = ["DeckBuilder", "Palette", "generation_stamp"]


def _left_run(text_frame):
    """Return a run on a left-aligned first paragraph of ``text_frame``.

    python-pptx text boxes inherit the presentation default alignment, which
    centres text. Every body paragraph in this deck is left-aligned, so this
    wrapper avoids repeating the two-line dance at each call site.

    Args:
        text_frame: The text frame to write into.

    Returns:
        A fresh run on the frame's first paragraph.

    Example:
        >>> from pptx import Presentation
        >>> slide = Presentation().slides.add_slide(
        ...     Presentation().slide_layouts[6])  # doctest: +SKIP
    """
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    return paragraph.add_run()


def generation_stamp(moment: datetime | None = None) -> str:
    """Render the generation time as a human-readable local timestamp.

    A date alone is not enough on a board that changes hourly: two decks
    generated the same day can disagree, and without a time nobody can tell
    which is current. The timezone is included because these decks travel
    between people in different ones.

    Args:
        moment: Override for the current time, for testing. Defaults to now
            in the local timezone.

    Returns:
        A string such as ``"23 August 2026 at 14:32 BST"``.

    Example:
        >>> from datetime import datetime, timezone
        >>> generation_stamp(datetime(2026, 8, 23, 14, 32, tzinfo=timezone.utc))
        '23 August 2026 at 14:32 UTC'
    """
    when = moment or datetime.now().astimezone()
    zone = when.strftime("%Z") or "local time"
    # The day comes from the integer rather than a strftime directive. The
    # no-padding directives are platform-specific — glibc and Windows spell
    # them differently, and each rejects the other's — so neither is portable.
    return f"{when.day} {when:%B %Y} at {when:%H:%M} {zone}"


class Palette:
    """Colours used throughout the deck.

    Deliberately not default-blue. Values are :class:`RGBColor` instances ready
    to assign to python-pptx colour properties.

    Attributes:
        ink: Primary text colour.
        muted: Secondary text colour.
        accent: Primary series and emphasis colour.
        accent_deep: Darker accent, for a second series that must not be
            mistaken for the first.
        accent_soft: Secondary series colour.
        warn: Colour for figures needing attention.
        rule: Hairline and gridline colour.
        surface: Card background tint.

    Example:
        >>> isinstance(Palette.accent, RGBColor)
        True
    """

    ink = RGBColor(0x1F, 0x29, 0x33)
    muted = RGBColor(0x61, 0x6E, 0x7C)
    accent = RGBColor(0x0F, 0x76, 0x6E)
    accent_deep = RGBColor(0x11, 0x5E, 0x59)
    accent_soft = RGBColor(0x7C, 0xC4, 0xBD)
    warn = RGBColor(0xB4, 0x53, 0x09)
    rule = RGBColor(0xDC, 0xE1, 0xE6)
    surface = RGBColor(0xF4, 0xF6, 0xF7)
    white = RGBColor(0xFF, 0xFF, 0xFF)


#: Slide canvas width for a 16:9 deck.
SLIDE_WIDTH = Inches(13.333)
#: Slide canvas height for a 16:9 deck.
SLIDE_HEIGHT = Inches(7.5)
#: Left and right margin for slide content.
MARGIN = Inches(0.7)


class DeckBuilder:
    """Assembles the sprint review deck.

    Args:
        project_title: Board or programme name shown on the title slide.
        subtitle: Optional line beneath the title, e.g. the entity name.

    Example:
        >>> DeckBuilder("Acme Platform").project_title
        'Acme Platform'
    """

    def __init__(
        self,
        project_title: str,
        subtitle: str = "",
        mode: str = "review",
    ) -> None:
        """Initialise the builder.

        Args:
            project_title: Board or programme name for the title slide.
            subtitle: Optional line beneath the title.
            mode: ``"review"`` for an end-of-sprint deck, or ``"midsprint"``
                for a check-in partway through. Mid-sprint decks describe
                incomplete work as in flight rather than as carryover, since
                nothing has rolled over yet.

        Raises:
            ValueError: If ``mode`` is not a recognised value.
        """
        if mode not in ("review", "midsprint"):
            raise ValueError(f"Unknown mode: {mode!r}")
        self.project_title = project_title
        self.subtitle = subtitle
        self.mode = mode
        self._generated_at = datetime.now().astimezone()

    def build(
        self,
        current: SprintMetrics,
        history: Sequence[SprintMetrics],
        burndown_points: Sequence[BurndownPoint],
        carryover: Sequence[ProjectItem],
        output_path: Path,
        milestone_forecasts: Sequence[tuple[str, float]] = (),
        unestimated: Sequence[ProjectItem] = (),
        unsprinted: Sequence[ProjectItem] = (),
        burndown_reconstructed: bool = False,
        churn: Sequence[ScopeChange] = (),
    ) -> Path:
        """Generate the deck and write it to disk.

        Args:
            current: Metrics for the sprint under review.
            history: Metrics for prior sprints, oldest first.
            burndown_points: Daily burndown curve for the current sprint.
            carryover: Incomplete items that will roll forward.
            output_path: Destination ``.pptx`` path.
            milestone_forecasts: ``(milestone, remaining_points)`` pairs to
                project completion dates for.
            unestimated: Items in the sprint carrying no estimate, listed with
                links so they can be corrected.
            unsprinted: Board items assigned to no iteration at all.
            churn: Day-by-day scope changes, from snapshot diffs.
            burndown_reconstructed: Whether the curve was derived from issue
                closure dates rather than daily snapshots. Annotated on the
                slide, since it cannot show mid-sprint scope changes.

        Trend slides (velocity, predictability, work mix) are omitted when
        ``history`` is empty, and the burndown is omitted when no snapshots
        cover the sprint, and the forecast when no milestone carries
        outstanding work. A single-column trend chart carries no information
        and reads as a defect; leaving the slide out beats shipping one bar
        floating in an empty plot area.

        Returns:
            The path written.

        Raises:
            OSError: If the file cannot be written.

        Example:
            >>> DeckBuilder("X").project_title
            'X'
        """
        presentation = Presentation()
        presentation.slide_width = SLIDE_WIDTH
        presentation.slide_height = SLIDE_HEIGHT

        self._title_slide(presentation, current)
        self._summary_slide(presentation, current, history)
        if burndown_points:
            self._burndown_slide(
                presentation, current, burndown_points, burndown_reconstructed
            )
            self._burnup_slide(
                presentation, current, burndown_points, burndown_reconstructed
            )
        if history:
            self._trend_slide(presentation, history, current)
        # Work mix needs no history: the planned/unplanned/carryover split of
        # a single sprint is a composition, not a trend, and is worth showing
        # from the first sprint onward.
        if current.committed_points > 0:
            self._work_mix_slide(presentation, history, current)
        if churn:
            self._scope_churn_slide(presentation, current, churn)
        self._carryover_slide(presentation, carryover)
        if milestone_forecasts:
            self._forecast_slide(
                presentation, history, current, milestone_forecasts
            )
        self._linked_table_slide(
            presentation,
            "Unestimated items",
            f"{len(unestimated)} item(s) in this sprint carry no estimate and "
            "count as zero in every figure above",
            unestimated,
            "Every item in this sprint carries an estimate.",
        )
        self._linked_table_slide(
            presentation,
            "Not in a sprint",
            f"{len(unsprinted)} board item(s) sit outside every iteration and "
            "appear in no report",
            unsprinted,
            "Every board item is assigned to an iteration.",
        )

        self._stamp_slides(presentation, current)
        self._stamp_metadata(presentation, current)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))
        return output_path

    def _stamp_slides(self, presentation: Presentation, current: SprintMetrics) -> None:
        """Add a generation footer to every slide except the title.

        Slides get screenshotted, pasted, and printed individually, so the
        stamp has to travel with each one rather than living only on slide 1.

        Args:
            presentation: The presentation being built.
            current: Metrics for the sprint, for the iteration name.
        """
        text = f"{current.iteration} · generated {generation_stamp(self._generated_at)}"
        for index, slide in enumerate(presentation.slides):
            if index == 0:
                continue
            box = slide.shapes.add_textbox(
                MARGIN, Inches(6.95), SLIDE_WIDTH - 2 * MARGIN, Inches(0.35)
            )
            frame = box.text_frame
            frame.word_wrap = True
            frame.margin_left = 0
            frame.margin_top = 0
            run = _left_run(frame)
            run.text = text
            run.font.size = Pt(9)
            run.font.color.rgb = Palette.muted

    def _stamp_metadata(
        self, presentation: Presentation, current: SprintMetrics
    ) -> None:
        """Write generation details into the file's own properties.

        Recorded in the document metadata as well as on the slides, so the
        provenance survives someone copying slides into another deck.

        Args:
            presentation: The presentation being built.
            current: Metrics for the sprint under review.
        """
        label = "Mid-Sprint Check" if self.mode == "midsprint" else "Sprint Review"
        properties = presentation.core_properties
        properties.title = f"{self.project_title} — {label} — {current.iteration}"
        properties.author = "gh-sprint-report"
        properties.comments = (
            f"Generated {generation_stamp(self._generated_at)} by "
            f"gh-sprint-report. Board data as at that moment."
        )
        properties.created = self._generated_at.replace(tzinfo=None)
        properties.modified = self._generated_at.replace(tzinfo=None)

    # ------------------------------------------------------------------
    # Slide-level helpers
    # ------------------------------------------------------------------

    def _blank(self, presentation: Presentation):
        """Add a blank slide.

        Args:
            presentation: The presentation being built.

        Returns:
            The newly added slide.
        """
        return presentation.slides.add_slide(presentation.slide_layouts[6])

    def _heading(self, slide, text: str, sub: str = "") -> None:
        """Draw a slide heading and optional sub-heading.

        Args:
            slide: Target slide.
            text: Heading text.
            sub: Optional smaller line beneath the heading.
        """
        box = slide.shapes.add_textbox(
            MARGIN, Inches(0.45), SLIDE_WIDTH - 2 * MARGIN, Inches(0.75)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_top = 0
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = Palette.ink

        if sub:
            sub_box = slide.shapes.add_textbox(
                MARGIN, Inches(1.12), SLIDE_WIDTH - 2 * MARGIN, Inches(0.4)
            )
            sub_frame = sub_box.text_frame
            sub_frame.word_wrap = True
            sub_frame.margin_left = 0
            sub_paragraph = sub_frame.paragraphs[0]
            sub_paragraph.alignment = PP_ALIGN.LEFT
            sub_run = sub_paragraph.add_run()
            sub_run.text = sub
            sub_run.font.size = Pt(14)
            sub_run.font.color.rgb = Palette.muted

    def _stat_card(
        self,
        slide,
        left: Emu,
        top: Emu,
        width: Emu,
        value: str,
        label: str,
        tone: RGBColor | None = None,
    ) -> None:
        """Draw a single headline figure on a tinted card.

        Args:
            slide: Target slide.
            left: Card x position.
            top: Card y position.
            width: Card width.
            value: The large figure.
            label: Caption beneath the figure.
            tone: Colour for the figure; defaults to the accent colour.
        """
        card = slide.shapes.add_shape(1, left, top, width, Inches(1.75))
        card.fill.solid()
        card.fill.fore_color.rgb = Palette.surface
        card.line.color.rgb = Palette.rule
        card.line.width = Pt(0.75)
        card.shadow.inherit = False

        frame = card.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.25)
        frame.margin_right = Inches(0.25)
        frame.margin_top = Inches(0.22)

        first = frame.paragraphs[0]
        first.alignment = PP_ALIGN.LEFT
        value_run = first.add_run()
        value_run.text = value
        value_run.font.size = Pt(40)
        value_run.font.bold = True
        value_run.font.color.rgb = tone or Palette.accent

        caption = frame.add_paragraph()
        caption.alignment = PP_ALIGN.LEFT
        caption_run = caption.add_run()
        caption_run.text = label
        caption_run.font.size = Pt(12)
        caption_run.font.color.rgb = Palette.muted

    def _style_chart(
        self,
        chart,
        show_legend: bool = False,
        hide_value_axis: bool = False,
    ) -> None:
        """Strip default chart furniture and apply the deck's typography.

        PowerPoint's defaults — heavy gridlines, axis lines, tick marks, and a
        value axis duplicating the data labels — read as a spreadsheet screen
        grab. Removing them is most of the difference between a chart that
        looks generated and one that looks designed.

        Args:
            chart: The chart object returned by ``add_chart``.
            show_legend: Whether to display the legend.
            hide_value_axis: Drop the value axis and gridlines entirely. Use
                when data labels already carry the numbers.
        """
        chart.has_title = False
        chart.font.size = Pt(12)
        chart.font.color.rgb = Palette.muted

        chart.has_legend = show_legend
        if show_legend:
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(12)
            chart.legend.font.color.rgb = Palette.muted

        value_axis = chart.value_axis
        value_axis.has_major_gridlines = not hide_value_axis
        if not hide_value_axis:
            gridlines = value_axis.major_gridlines.format.line
            gridlines.color.rgb = Palette.rule
            gridlines.width = Pt(0.75)
        value_axis.visible = not hide_value_axis
        value_axis.has_minor_gridlines = False
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.minor_tick_mark = XL_TICK_MARK.NONE
        value_axis.format.line.fill.background()
        value_axis.tick_labels.font.size = Pt(11)
        value_axis.tick_labels.font.color.rgb = Palette.muted

        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.minor_tick_mark = XL_TICK_MARK.NONE
        category_axis.format.line.color.rgb = Palette.rule
        category_axis.format.line.width = Pt(0.75)
        category_axis.tick_labels.font.size = Pt(13)
        category_axis.tick_labels.font.color.rgb = Palette.ink

    # ------------------------------------------------------------------
    # Slides
    # ------------------------------------------------------------------

    def _title_slide(self, presentation: Presentation, current: SprintMetrics) -> None:
        """Build the title slide.

        Args:
            presentation: The presentation being built.
            current: Metrics for the sprint under review.
        """
        slide = self._blank(presentation)
        box = slide.shapes.add_textbox(
            MARGIN, Inches(2.6), SLIDE_WIDTH - 2 * MARGIN, Inches(2.2)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        title = frame.paragraphs[0].add_run()
        label = "Mid-Sprint Check" if self.mode == "midsprint" else "Sprint Review"
        title.text = f"{self.project_title} — {label}"
        title.font.size = Pt(44)
        title.font.bold = True
        title.font.color.rgb = Palette.ink

        detail = frame.add_paragraph()
        detail.alignment = PP_ALIGN.LEFT
        detail.space_before = Pt(10)
        window = ""
        if current.start and current.end:
            window = f"  ·  {current.start:%d %b} – {current.end:%d %b %Y}"
        detail_run = detail.add_run()
        detail_run.text = f"{current.iteration}{window}"
        detail_run.font.size = Pt(20)
        detail_run.font.color.rgb = Palette.accent

        if self.subtitle:
            sub = frame.add_paragraph()
            sub.alignment = PP_ALIGN.LEFT
            sub.space_before = Pt(14)
            sub_run = sub.add_run()
            sub_run.text = self.subtitle
            sub_run.font.size = Pt(14)
            sub_run.font.color.rgb = Palette.muted

        footer = slide.shapes.add_textbox(
            MARGIN, Inches(6.6), SLIDE_WIDTH - 2 * MARGIN, Inches(0.4)
        ).text_frame
        footer.word_wrap = True
        footer.margin_left = 0
        footer_paragraph = footer.paragraphs[0]
        footer_paragraph.alignment = PP_ALIGN.LEFT
        footer_run = footer_paragraph.add_run()
        footer_run.text = f"Generated {generation_stamp(self._generated_at)}"
        footer_run.font.size = Pt(12)
        footer_run.font.color.rgb = Palette.muted

    def _summary_slide(
        self,
        presentation: Presentation,
        current: SprintMetrics,
        history: Sequence[SprintMetrics],
    ) -> None:
        """Build the headline-figures slide.

        Args:
            presentation: The presentation being built.
            current: Metrics for the sprint under review.
            history: Prior sprint metrics for the rolling average.
        """
        slide = self._blank(presentation)
        self._heading(
            slide,
            "Sprint at a glance",
            f"{current.completed_items} of {current.total_items} items complete",
        )

        gutter = Inches(0.3)
        card_width = Emu(int((SLIDE_WIDTH - 2 * MARGIN - 3 * gutter) / 4))
        average = rolling_average(
            [metric.completed_points for metric in history] or [0.0]
        )

        cards = [
            (f"{current.completed_points:g}", "Points completed", Palette.accent),
            (f"{current.predictability:g}%", "Of commitment delivered", Palette.accent),
            (
                f"{current.unplanned_share:g}%",
                "Unplanned work",
                Palette.warn if current.unplanned_share > 20 else Palette.accent,
            ),
            (f"{average:g}", "3-sprint average", Palette.muted),
        ]
        for index, (value, label, tone) in enumerate(cards):
            left = Emu(int(MARGIN + index * (card_width + gutter)))
            self._stat_card(slide, left, Inches(1.9), card_width, value, label, tone)

        note_lines = [
            f"Committed {current.committed_points:g} pts across "
            f"{current.total_items} items.",
            f"Planned {current.planned_points:g} · "
            f"Unplanned {current.unplanned_points:g} · "
            f"Carried in {current.carryover_points:g}.",
        ]
        if current.unestimated_items:
            note_lines.append(
                f"{current.unestimated_items} item(s) carry no estimate and "
                "count as zero."
            )

        box = slide.shapes.add_textbox(
            MARGIN, Inches(4.2), SLIDE_WIDTH - 2 * MARGIN, Inches(2.0)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        for index, line in enumerate(note_lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)
            run = paragraph.add_run()
            run.text = line
            run.font.size = Pt(15)
            run.font.color.rgb = Palette.ink

    def _burndown_slide(
        self,
        presentation: Presentation,
        current: SprintMetrics,
        points: Sequence[BurndownPoint],
        reconstructed: bool = False,
    ) -> None:
        """Build the burndown slide.

        Args:
            presentation: The presentation being built.
            current: Metrics for the sprint under review.
            points: The burndown curve; an empty sequence renders a notice.
            reconstructed: Whether the curve came from closure dates rather
                than snapshots.
        """
        slide = self._blank(presentation)
        self._heading(slide, "Burndown", f"Remaining points across {current.iteration}")

        if not points:
            self._empty_notice(
                slide,
                "No snapshots recorded for this sprint yet. Burndown appears "
                "once the daily collector has run during a sprint.",
            )
            return

        chart_data = CategoryChartData()
        chart_data.categories = [point.day.strftime("%d %b") for point in points]
        chart_data.add_series("Remaining", [point.remaining for point in points])
        chart_data.add_series("Ideal", [point.ideal for point in points])

        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            MARGIN,
            Inches(1.75),
            SLIDE_WIDTH - 2 * MARGIN,
            Inches(5.0),
            chart_data,
        )
        chart = frame.chart
        self._style_chart(chart, show_legend=True)
        chart.series[0].format.line.color.rgb = Palette.accent
        chart.series[0].format.line.width = Pt(2.75)
        chart.series[1].format.line.color.rgb = Palette.rule
        chart.series[1].format.line.width = Pt(1.5)

    def _burnup_slide(
        self,
        presentation: Presentation,
        current: SprintMetrics,
        points: Sequence[BurndownPoint],
        reconstructed: bool = False,
    ) -> None:
        """Build the burnup slide.

        The same daily data as the burndown, read the other way: completed
        work rising toward a scope line. This is the framing that makes scope
        change visible — a burndown flattening could mean work stalled or
        work being added, and only a burnup distinguishes them.

        Args:
            presentation: The presentation being built.
            current: Metrics for the sprint under review.
            points: The daily curve.
            reconstructed: Whether the curve came from closure dates, in which
                case scope is flat because nothing records when items joined.
        """
        slide = self._blank(presentation)
        self._heading(
            slide,
            "Burnup",
            "Completed work rising toward total scope"
            + (
                " · scope is flat here because closure dates cannot say when "
                "an item joined the sprint"
                if reconstructed
                else " · a rising scope line is work added mid-sprint"
            ),
        )

        chart_data = CategoryChartData()
        chart_data.categories = [point.day.strftime("%d %b") for point in points]
        chart_data.add_series("Completed", [point.completed for point in points])
        chart_data.add_series("Scope", [point.scope for point in points])

        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            MARGIN,
            Inches(1.9),
            SLIDE_WIDTH - 2 * MARGIN,
            Inches(4.7),
            chart_data,
        )
        chart = frame.chart
        self._style_chart(chart, show_legend=True)
        chart.series[0].format.line.color.rgb = Palette.accent
        chart.series[0].format.line.width = Pt(2.75)
        chart.series[1].format.line.color.rgb = Palette.rule
        chart.series[1].format.line.width = Pt(2)

    def _scope_churn_slide(
        self,
        presentation: Presentation,
        current: SprintMetrics,
        changes: Sequence[ScopeChange],
    ) -> None:
        """Build the scope churn slide.

        What entered and left the sprint after it started. GitHub keeps no
        history of iteration-field changes, so this comes from diffing daily
        snapshots — which means it is only as granular as the collector's
        schedule, and only exists from the day it was switched on.

        This is usually the most useful slide for a leadership audience: it
        answers "did the plan hold" with specifics rather than a percentage.

        Args:
            presentation: The presentation being built.
            current: Metrics for the sprint under review.
            changes: Day-by-day membership changes.
        """
        slide = self._blank(presentation)
        added = sum(change.added_points for change in changes)
        removed = sum(change.removed_points for change in changes)
        net = added - removed
        direction = "grew" if net > 0 else "shrank" if net < 0 else "held"
        self._heading(
            slide,
            "Scope churn",
            f"Scope {direction} by {abs(net):g} points after the sprint started"
            if net
            else "Scope held from planning to close",
        )

        if not changes:
            self._empty_notice(
                slide,
                "No scope changes recorded. Either nothing moved, or daily "
                "snapshots did not cover this sprint.",
            )
            return

        rows = [
            (change.day, item, "In")
            for change in changes
            for item in change.added
        ] + [
            (change.day, item, "Out")
            for change in changes
            for item in change.removed
        ]
        rows.sort(key=lambda row: (row[0], row[2]))
        shown = rows[:10]

        table = slide.shapes.add_table(
            len(shown) + 1,
            4,
            MARGIN,
            Inches(1.9),
            SLIDE_WIDTH - 2 * MARGIN,
            Inches(0.34) * (len(shown) + 1),
        ).table
        table.first_row = True
        table.horz_banding = False
        table.columns[0].width = Inches(1.2)
        table.columns[1].width = Inches(1.0)
        table.columns[2].width = Inches(8.7)
        table.columns[3].width = Inches(1.0)

        for column, label in enumerate(("Day", "In/Out", "Item", "Points")):
            cell = table.cell(0, column)
            cell.text = label
            self._style_cell(cell, bold=True, colour=Palette.white, size=Pt(12))
            cell.fill.solid()
            cell.fill.fore_color.rgb = Palette.ink

        for row, (day, item, direction_label) in enumerate(shown, start=1):
            table.cell(row, 0).text = f"{day:%d %b}"
            table.cell(row, 1).text = direction_label
            title = item.title if len(item.title) <= 74 else item.title[:71] + "..."
            table.cell(row, 2).text = title
            table.cell(row, 3).text = f"{item.effective_points:g}"
            for column in range(4):
                cell = table.cell(row, column)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    Palette.surface if row % 2 else Palette.white
                )
                self._style_cell(cell, size=Pt(12))

        summary = f"{added:g} points added, {removed:g} removed"
        if len(rows) > len(shown):
            summary += f" · {len(rows) - len(shown)} further change(s) not shown"
        box = slide.shapes.add_textbox(
            MARGIN, Inches(6.4), SLIDE_WIDTH - 2 * MARGIN, Inches(0.5)
        )
        box.text_frame.word_wrap = True
        box.text_frame.margin_left = 0
        run = _left_run(box.text_frame)
        run.text = summary
        run.font.size = Pt(15)
        run.font.color.rgb = Palette.warn if added > removed else Palette.ink

    def _trend_slide(
        self,
        presentation: Presentation,
        history: Sequence[SprintMetrics],
        current: SprintMetrics,
    ) -> None:
        """Build the velocity slide.

        Committed and completed points per sprint, drawn as lines because the
        question is direction over time rather than the size of any one
        sprint. The gap between the lines is what was not delivered, and the
        percentages beneath give it as a ratio.

        Args:
            presentation: The presentation being built.
            history: Prior sprint metrics, oldest first.
            current: Metrics for the sprint under review.
        """
        slide = self._blank(presentation)
        series = list(history) + [current]
        average = rolling_average([metric.completed_points for metric in series])
        self._heading(
            slide,
            "Velocity",
            f"Points per sprint · recent average {average:g} completed",
        )

        chart_data = CategoryChartData()
        chart_data.categories = [metric.iteration for metric in series]
        chart_data.add_series(
            "Completed", [metric.completed_points for metric in series]
        )
        chart_data.add_series(
            "Committed", [metric.committed_points for metric in series]
        )
        chart_data.add_series("Average", [average] * len(series))

        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            MARGIN,
            Inches(1.9),
            SLIDE_WIDTH - 2 * MARGIN,
            Inches(4.1),
            chart_data,
        )
        chart = frame.chart
        self._style_chart(chart, show_legend=True)
        chart.series[0].format.line.color.rgb = Palette.accent
        chart.series[0].format.line.width = Pt(2.75)
        chart.series[1].format.line.color.rgb = Palette.accent_soft
        chart.series[1].format.line.width = Pt(2)
        chart.series[2].format.line.color.rgb = Palette.rule
        chart.series[2].format.line.width = Pt(1.25)

        ratios = "    ".join(
            f"{metric.iteration}  {metric.predictability:g}%"
            for metric in series[-6:]
        )
        box = slide.shapes.add_textbox(
            MARGIN, Inches(6.2), SLIDE_WIDTH - 2 * MARGIN, Inches(0.8)
        )
        box.text_frame.word_wrap = True
        box.text_frame.margin_left = 0
        label = _left_run(box.text_frame)
        label.text = "Delivered against commitment"
        label.font.size = Pt(12)
        label.font.color.rgb = Palette.muted

        values = box.text_frame.add_paragraph()
        values.alignment = PP_ALIGN.LEFT
        values_run = values.add_run()
        values_run.text = ratios
        values_run.font.size = Pt(16)
        values_run.font.color.rgb = Palette.ink

    def _work_mix_slide(
        self,
        presentation: Presentation,
        history: Sequence[SprintMetrics],
        current: SprintMetrics,
    ) -> None:
        """Build the work mix slide, showing both directions of the sprint.

        Two bars that sum to the same total. **Came from** splits the
        committed points by origin — planned, pulled in mid-sprint, or rolled
        in from the sprint before. **Went to** splits the same points by
        outcome — completed, or rolling forward to the next sprint.

        Showing only one direction invites the reader to supply the other from
        assumption, and the word "carryover" means opposite things depending
        on which they assume.

        Args:
            presentation: The presentation being built.
            history: Prior sprint metrics; used only to tell whether a
                previous sprint existed at all.
            current: Metrics for the sprint under review.
        """
        slide = self._blank(presentation)
        self._heading(
            slide,
            "Work mix",
            "Where this sprint's work came from, and where it ended up",
        )

        chart_data = CategoryChartData()
        chart_data.categories = ["Came from", "Went to"]
        # Series order is deliberately chronological, so both bars read left
        # to right as a timeline: work already in hand, work committed at
        # planning, work that arrived later — then what closed and what leaves.
        chart_data.add_series("Carried in", (current.carryover_points, 0))
        chart_data.add_series("Planned", (current.planned_points, 0))
        chart_data.add_series("Unplanned", (current.unplanned_points, 0))
        chart_data.add_series("Completed", (0, current.completed_points))
        chart_data.add_series("Rolling forward", (0, current.remaining_points))

        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            MARGIN,
            Inches(1.95),
            SLIDE_WIDTH - 2 * MARGIN,
            Inches(2.9),
            chart_data,
        )
        chart = frame.chart
        self._style_chart(chart, show_legend=True, hide_value_axis=True)

        colours = (
            Palette.accent_soft,
            Palette.accent,
            Palette.warn,
            Palette.accent_deep,
            Palette.rule,
        )
        for index, colour in enumerate(colours):
            chart.series[index].format.fill.solid()
            chart.series[index].format.fill.fore_color.rgb = colour

        plot = chart.plots[0]
        plot.gap_width = 80
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(12)
        plot.data_labels.font.bold = True
        plot.data_labels.font.color.rgb = Palette.white
        # Each series is present in only one of the two bars, so the other
        # holds a zero. A number format whose zero section is empty hides
        # those labels; setting them per point is not exposed by python-pptx.
        plot.data_labels.number_format = "0;;;"
        plot.data_labels.number_format_is_linked = False

        lines = [
            f"Interrupt load: {current.unplanned_share:g}% of committed points "
            "arrived after planning.",
        ]
        tone = Palette.warn if current.unplanned_share > 20 else Palette.ink

        # A first sprint has nothing behind it, so carry-in is impossible.
        # Say so rather than charting a figure that cannot be true.
        if not history and current.carryover_points > 0:
            lines.append(
                f"{current.carryover_points:g} points are marked "
                "'Carried in', but this is the first sprint — nothing existed "
                "to carry in from. Those items are most likely mislabelled; "
                "work leaving for the next sprint is counted under "
                "'Rolling forward'."
            )
            tone = Palette.warn

        box = slide.shapes.add_textbox(
            MARGIN, Inches(5.1), SLIDE_WIDTH - 2 * MARGIN, Inches(1.6)
        )
        box.text_frame.word_wrap = True
        box.text_frame.margin_left = 0
        for index, line in enumerate(lines):
            paragraph = (
                box.text_frame.paragraphs[0]
                if index == 0
                else box.text_frame.add_paragraph()
            )
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)
            run = paragraph.add_run()
            run.text = line
            run.font.size = Pt(15) if index == 0 else Pt(14)
            run.font.color.rgb = tone if index == 0 else Palette.warn

    def _carryover_slide(
        self, presentation: Presentation, carryover: Sequence[ProjectItem]
    ) -> None:
        """Build the carryover detail slide.

        Args:
            presentation: The presentation being built.
            carryover: Incomplete items rolling to the next sprint.
        """
        slide = self._blank(presentation)
        total = sum(item.effective_points for item in carryover)
        if self.mode == "midsprint":
            heading, sub, empty = (
                "In flight",
                f"{len(carryover)} item(s), {total:g} points still open",
                "Everything committed is already closed.",
            )
        else:
            heading, sub, empty = (
                "Rolling forward",
                f"{len(carryover)} item(s), {total:g} points moving to the "
                "next sprint",
                "Nothing rolls forward. Everything committed closed.",
            )
        self._heading(slide, heading, sub)

        if not carryover:
            self._empty_notice(slide, empty)
            return

        rows = min(len(carryover), 8) + 1
        table_shape = slide.shapes.add_table(
            rows,
            3,
            MARGIN,
            Inches(1.75),
            SLIDE_WIDTH - 2 * MARGIN,
            Inches(0.45) * rows,
        )
        table = table_shape.table
        table.first_row = True
        table.horz_banding = False
        table.columns[0].width = Inches(8.0)
        table.columns[1].width = Inches(1.9)
        table.columns[2].width = Inches(2.0)

        for column, label in enumerate(("Item", "Points", "Status")):
            cell = table.cell(0, column)
            cell.text = label
            self._style_cell(cell, bold=True, colour=Palette.white)
            cell.fill.solid()
            cell.fill.fore_color.rgb = Palette.ink

        for row, item in enumerate(carryover[:8], start=1):
            title = item.title if len(item.title) <= 78 else item.title[:75] + "…"
            table.cell(row, 0).text = title
            table.cell(row, 1).text = f"{item.effective_points:g}"
            table.cell(row, 2).text = item.status or "—"
            for column in range(3):
                cell = table.cell(row, column)
                cell.fill.solid()
                cell.fill.fore_color.rgb = Palette.surface if row % 2 else Palette.white
                self._style_cell(cell)

        if len(carryover) > 8:
            box = slide.shapes.add_textbox(
                MARGIN, Inches(6.6), SLIDE_WIDTH - 2 * MARGIN, Inches(0.4)
            )
            box.text_frame.word_wrap = True
            box.text_frame.margin_left = 0
            run = _left_run(box.text_frame)
            run.text = f"+ {len(carryover) - 8} further item(s) not shown."
            run.font.size = Pt(12)
            run.font.color.rgb = Palette.muted

    def _linked_table_slide(
        self,
        presentation: Presentation,
        heading: str,
        sub: str,
        items: Sequence[ProjectItem],
        empty_message: str,
        max_rows: int = 13,
    ) -> None:
        """Render a slide listing items as a table of clickable links.

        Used for the follow-up lists — unestimated items and items sitting
        outside any sprint — where the point of the slide is that someone can
        click straight through and fix the row.

        Args:
            presentation: The presentation being built.
            heading: Slide heading.
            sub: Sub-heading beneath it.
            items: Items to list; each contributes one row.
            empty_message: Text shown when ``items`` is empty.
            max_rows: Rows to render before summarising the remainder.
        """
        slide = self._blank(presentation)
        self._heading(slide, heading, sub)

        if not items:
            self._empty_notice(slide, empty_message)
            return

        shown = list(items)[:max_rows]
        table = slide.shapes.add_table(
            len(shown) + 1,
            3,
            MARGIN,
            Inches(1.75),
            SLIDE_WIDTH - 2 * MARGIN,
            Inches(0.32) * (len(shown) + 1),
        ).table
        table.first_row = True
        table.horz_banding = False
        table.columns[0].width = Inches(0.9)
        table.columns[1].width = Inches(9.1)
        table.columns[2].width = Inches(1.9)

        for column, label in enumerate(("#", "Item", "Status")):
            cell = table.cell(0, column)
            cell.text = label
            self._style_cell(cell, bold=True, colour=Palette.white, size=Pt(12))
            cell.fill.solid()
            cell.fill.fore_color.rgb = Palette.ink

        for row, item in enumerate(shown, start=1):
            number = table.cell(row, 0)
            number.text = f"#{item.item_id}"

            link_cell = table.cell(row, 1)
            link_cell.text = ""
            paragraph = link_cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            title = item.title if len(item.title) <= 88 else item.title[:85] + "…"
            run.text = title
            if item.url:
                run.hyperlink.address = item.url

            status = table.cell(row, 2)
            status.text = item.status or "—"

            for column in range(3):
                cell = table.cell(row, column)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    Palette.surface if row % 2 else Palette.white
                )
                self._style_cell(cell, size=Pt(12))
            if item.url:
                run.font.color.rgb = Palette.accent
                run.font.underline = True

        if len(items) > max_rows:
            box = slide.shapes.add_textbox(
                MARGIN, Inches(6.7), SLIDE_WIDTH - 2 * MARGIN, Inches(0.4)
            )
            box.text_frame.word_wrap = True
            box.text_frame.margin_left = 0
            note = _left_run(box.text_frame)
            note.text = f"+ {len(items) - max_rows} further item(s) not shown."
            note.font.size = Pt(12)
            note.font.color.rgb = Palette.muted

    def _forecast_slide(
        self,
        presentation: Presentation,
        history: Sequence[SprintMetrics],
        current: SprintMetrics,
        milestone_forecasts: Sequence[tuple[str, float]],
    ) -> None:
        """Build the forecast slide.

        Args:
            presentation: The presentation being built.
            history: Prior sprint metrics, oldest first.
            current: Metrics for the sprint under review.
            milestone_forecasts: ``(milestone, remaining_points)`` pairs.
        """
        slide = self._blank(presentation)
        series = list(history) + [current]
        velocity = rolling_average([metric.completed_points for metric in series])
        self._heading(slide, "Forecast", f"Projected at {velocity:g} points per sprint")

        if not milestone_forecasts:
            self._empty_notice(
                slide,
                "No milestones with outstanding work. Assign a milestone to "
                "open items to project delivery dates.",
            )
            return

        rows = len(milestone_forecasts) + 1
        table = slide.shapes.add_table(
            rows,
            3,
            MARGIN,
            Inches(1.9),
            SLIDE_WIDTH - 2 * MARGIN,
            Inches(0.45) * rows,
        ).table
        table.first_row = True
        table.horz_banding = False
        table.columns[0].width = Inches(7.0)
        table.columns[1].width = Inches(2.4)
        table.columns[2].width = Inches(2.5)

        for column, label in enumerate(("Milestone", "Points remaining", "Sprints")):
            cell = table.cell(0, column)
            cell.text = label
            self._style_cell(cell, bold=True, colour=Palette.white)
            cell.fill.solid()
            cell.fill.fore_color.rgb = Palette.ink

        for row, (milestone, remaining) in enumerate(milestone_forecasts, start=1):
            sprints = forecast_sprints(remaining, velocity)
            table.cell(row, 0).text = milestone
            table.cell(row, 1).text = f"{remaining:g}"
            table.cell(row, 2).text = "—" if sprints is None else f"{sprints:g}"
            for column in range(3):
                cell = table.cell(row, column)
                cell.fill.solid()
                cell.fill.fore_color.rgb = Palette.surface if row % 2 else Palette.white
                self._style_cell(cell)

    # ------------------------------------------------------------------
    # Small shared pieces
    # ------------------------------------------------------------------

    def _style_cell(
        self,
        cell,
        bold: bool = False,
        colour: RGBColor | None = None,
        size: Pt | None = None,
    ) -> None:
        """Apply table cell typography.

        Args:
            cell: The table cell to style.
            bold: Whether to bold the text.
            colour: Text colour; defaults to the ink colour.
            size: Font size; defaults to 13pt.
        """
        cell.margin_left = Inches(0.12)
        cell.margin_top = Inches(0.04)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = size or Pt(13)
                run.font.bold = bold
                if run.font.color.type is None:
                    run.font.color.rgb = colour or Palette.ink

    def _empty_notice(self, slide, message: str) -> None:
        """Render a centred explanatory message on an otherwise empty slide.

        Args:
            slide: Target slide.
            message: Text to display.
        """
        box = slide.shapes.add_textbox(
            MARGIN, Inches(3.0), SLIDE_WIDTH - 2 * MARGIN, Inches(1.4)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        run = _left_run(frame)
        run.text = message
        run.font.size = Pt(17)
        run.font.color.rgb = Palette.muted
