"""Tests for scope-change detection and the churn slide."""

from __future__ import annotations

from datetime import date
from pathlib import Path

from sprint_report.deck import DeckBuilder
from sprint_report.metrics import ScopeChange, scope_changes
from sprint_report.models import ProjectItem, Snapshot, SprintMetrics


def make(ident, points):
    """Build a sprint item.

    Args:
        ident: Item identifier.
        points: Estimate.

    Returns:
        A project item in Sprint 1.
    """
    return ProjectItem(
        item_id=str(ident),
        title=f"Item {ident}",
        iteration="Sprint 1",
        iteration_start=date(2026, 8, 10),
        iteration_duration=14,
        points=points,
    )


def snap(day, items):
    """Build a snapshot.

    Args:
        day: Day of August 2026.
        items: Items on the board.

    Returns:
        A snapshot for that day.
    """
    return Snapshot(date(2026, 8, day), "Board", items)


class TestScopeChanges:
    """Diffing consecutive snapshots."""

    def test_detects_an_addition(self):
        """An item present today and absent yesterday came in."""
        changes = scope_changes(
            [snap(10, [make(1, 8)]), snap(11, [make(1, 8), make(2, 5)])],
            "Sprint 1",
        )
        assert len(changes) == 1
        assert [i.item_id for i in changes[0].added] == ["2"]
        assert changes[0].added_points == 5.0

    def test_detects_a_removal(self):
        """An item absent today and present yesterday went out."""
        changes = scope_changes(
            [snap(10, [make(1, 8), make(2, 5)]), snap(11, [make(1, 8)])],
            "Sprint 1",
        )
        assert [i.item_id for i in changes[0].removed] == ["2"]
        assert changes[0].removed_points == 5.0

    def test_net_points(self):
        """Net is what the sprint's committed total moved by."""
        changes = scope_changes(
            [
                snap(10, [make(1, 8), make(2, 5)]),
                snap(11, [make(1, 8), make(3, 13)]),
            ],
            "Sprint 1",
        )
        assert changes[0].net_points == 8.0

    def test_quiet_days_are_omitted(self):
        """Only days where membership moved produce a record."""
        board = [make(1, 8)]
        changes = scope_changes(
            [snap(10, board), snap(11, board), snap(12, board)], "Sprint 1"
        )
        assert changes == []

    def test_multiple_days(self):
        """Each changing day is reported separately."""
        changes = scope_changes(
            [
                snap(10, [make(1, 8)]),
                snap(11, [make(1, 8), make(2, 5)]),
                snap(12, [make(1, 8), make(2, 5), make(3, 3)]),
            ],
            "Sprint 1",
        )
        assert [c.day.day for c in changes] == [11, 12]

    def test_first_snapshot_is_the_baseline(self):
        """Opening membership is not itself a change."""
        changes = scope_changes([snap(10, [make(1, 8), make(2, 5)])], "Sprint 1")
        assert changes == []

    def test_needs_two_snapshots(self):
        """A single day cannot be diffed against anything."""
        assert scope_changes([snap(10, [make(1, 8)])], "Sprint 1") == []

    def test_no_snapshots(self):
        """Nothing to diff."""
        assert scope_changes([], "Sprint 1") == []

    def test_unordered_input_is_sorted(self):
        """Snapshots may arrive in any order."""
        changes = scope_changes(
            [snap(11, [make(1, 8), make(2, 5)]), snap(10, [make(1, 8)])],
            "Sprint 1",
        )
        assert changes[0].day == date(2026, 8, 11)
        assert [i.item_id for i in changes[0].added] == ["2"]

    def test_other_iterations_are_ignored(self):
        """A change in Sprint 2 is not a change in Sprint 1."""
        other = ProjectItem(
            item_id="9", title="t", iteration="Sprint 2", points=5,
            iteration_start=date(2026, 8, 24), iteration_duration=14,
        )
        changes = scope_changes(
            [snap(10, [make(1, 8)]), snap(11, [make(1, 8), other])], "Sprint 1"
        )
        assert changes == []


class TestChurnSlide:
    """Rendering of the churn slide."""

    def _build(self, tmp_path: Path, churn):
        """Generate a deck with churn data.

        Args:
            tmp_path: pytest temporary directory.
            churn: Scope changes to render.

        Returns:
            The generated deck path.
        """
        return DeckBuilder("X").build(
            current=SprintMetrics(
                "Sprint 1", committed_points=40, planned_points=40
            ),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "churn.pptx",
            churn=churn,
        )

    def _text(self, path: Path) -> str:
        """Flatten shape and table text.

        Args:
            path: Deck to read.

        Returns:
            All text joined by newlines.
        """
        from pptx import Presentation

        chunks = []
        for slide in Presentation(str(path)).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            chunks.append(cell.text)
        return "\n".join(chunks)

    def test_renders_additions_and_removals(self, tmp_path: Path):
        """Both directions appear in the table."""
        churn = [
            ScopeChange(date(2026, 8, 12), [make(2, 5)], [make(3, 3)]),
        ]
        text = self._text(self._build(tmp_path, churn))
        assert "Scope churn" in text
        assert "Item 2" in text
        assert "Item 3" in text
        assert "In" in text and "Out" in text

    def test_summarises_net_growth(self, tmp_path: Path):
        """The heading states the direction and size."""
        churn = [ScopeChange(date(2026, 8, 12), [make(2, 22)], [])]
        assert "grew by 22 points" in self._text(self._build(tmp_path, churn))

    def test_summarises_net_shrink(self, tmp_path: Path):
        """Work removed reads as a shrink, not a negative growth."""
        churn = [ScopeChange(date(2026, 8, 12), [], [make(2, 8)])]
        assert "shrank by 8 points" in self._text(self._build(tmp_path, churn))

    def test_no_slide_without_churn_data(self, tmp_path: Path):
        """Absent snapshots mean no slide at all, not an empty one."""
        assert "Scope churn" not in self._text(self._build(tmp_path, []))

    def test_long_lists_are_truncated(self, tmp_path: Path):
        """Beyond ten rows the remainder is summarised."""
        churn = [
            ScopeChange(date(2026, 8, 12), [make(i, 3) for i in range(15)], [])
        ]
        assert "further change(s) not shown" in self._text(
            self._build(tmp_path, churn)
        )
