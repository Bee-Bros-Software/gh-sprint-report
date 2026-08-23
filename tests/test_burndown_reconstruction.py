"""Tests for reconstructing a burndown from issue closure dates."""

from __future__ import annotations

from datetime import date

from sprint_report.metrics import burndown_from_closures
from sprint_report.models import ProjectItem

START = date(2026, 8, 10)
END = date(2026, 8, 23)


def item(ident, points, closed_at=None, status="In Progress"):
    """Build a sprint item.

    Args:
        ident: Item identifier.
        points: Estimate.
        closed_at: Closure date, if closed.
        status: Board status.

    Returns:
        A configured project item.
    """
    return ProjectItem(
        item_id=str(ident),
        title=f"Item {ident}",
        status="Done" if closed_at else status,
        iteration="Sprint 1",
        iteration_start=START,
        iteration_duration=14,
        points=points,
        closed=closed_at is not None,
        closed_at=closed_at,
    )


class TestReconstruction:
    """Curve shape and arithmetic."""

    def _board(self):
        """A sprint with 40 points, half of it closed on known dates.

        Returns:
            Five project items.
        """
        return [
            item(1, 8, date(2026, 8, 12)),
            item(2, 5, date(2026, 8, 14)),
            item(3, 7, date(2026, 8, 18)),
            item(4, 12),
            item(5, 8),
        ]

    def test_one_point_per_day(self):
        """The curve covers every day from start to today or sprint end."""
        curve = burndown_from_closures(
            self._board(), "Sprint 1", today=date(2026, 8, 16)
        )
        assert [p.day for p in curve] == [
            date(2026, 8, d) for d in range(10, 17)
        ]

    def test_starts_at_full_scope(self):
        """Day one has nothing burned down yet."""
        curve = burndown_from_closures(
            self._board(), "Sprint 1", today=date(2026, 8, 16)
        )
        assert curve[0].remaining == 40.0

    def test_remaining_drops_on_closure_days(self):
        """Closing 8 points on the 12th moves the line by 8."""
        curve = {
            p.day: p.remaining
            for p in burndown_from_closures(
                self._board(), "Sprint 1", today=date(2026, 8, 16)
            )
        }
        assert curve[date(2026, 8, 11)] == 40.0
        assert curve[date(2026, 8, 12)] == 32.0
        assert curve[date(2026, 8, 13)] == 32.0
        assert curve[date(2026, 8, 14)] == 27.0

    def test_stops_at_today_mid_sprint(self):
        """No line is drawn into the future."""
        curve = burndown_from_closures(
            self._board(), "Sprint 1", today=date(2026, 8, 16)
        )
        assert curve[-1].day == date(2026, 8, 16)

    def test_runs_to_sprint_end_when_closed(self):
        """A finished sprint charts its whole span."""
        curve = burndown_from_closures(
            self._board(), "Sprint 1", today=date(2026, 9, 1)
        )
        assert curve[-1].day == END

    def test_ideal_line_falls_to_zero(self):
        """The ideal runs linearly from scope to nothing."""
        curve = burndown_from_closures(
            self._board(), "Sprint 1", today=date(2026, 9, 1)
        )
        assert curve[0].ideal == 40.0
        assert curve[-1].ideal == 0.0

    def test_done_without_a_closure_date(self):
        """Board-status completion has no timestamp; count it on the last day."""
        board = [item(1, 10, status="Done"), item(2, 10)]
        board[0] = ProjectItem(
            item_id="1", title="t", status="Done", iteration="Sprint 1",
            iteration_start=START, iteration_duration=14, points=10,
            closed=False, closed_at=None,
        )
        curve = burndown_from_closures(board, "Sprint 1", today=date(2026, 8, 14))
        assert curve[0].remaining == 20.0
        assert curve[-1].remaining == 10.0

    def test_closure_before_sprint_start_is_clamped(self):
        """An item closed before the sprint began burns on day one."""
        board = [item(1, 5, date(2026, 7, 1)), item(2, 5)]
        curve = burndown_from_closures(board, "Sprint 1", today=date(2026, 8, 12))
        assert curve[0].remaining == 5.0


class TestEdgeCases:
    """Conditions where no curve can be produced."""

    def test_unknown_iteration(self):
        """Nothing to chart."""
        assert burndown_from_closures([], "Sprint 9") == []

    def test_iteration_without_dates(self):
        """No timeline means no curve."""
        board = [ProjectItem(item_id="1", title="t", iteration="S1", points=5)]
        assert burndown_from_closures(board, "S1") == []

    def test_no_estimated_work(self):
        """A sprint of unestimated items has no points to burn."""
        board = [item(1, None), item(2, None)]
        assert burndown_from_closures(board, "Sprint 1", today=END) == []

    def test_sprint_entirely_in_the_future(self):
        """A sprint that has not started yet charts nothing."""
        board = [item(1, 5)]
        assert burndown_from_closures(
            board, "Sprint 1", today=date(2026, 8, 1)
        ) == []


class TestScopeLimitation:
    """The documented inaccuracy is real and worth pinning down."""

    def test_late_additions_counted_from_day_one(self):
        """Nothing records when an item joined, so scope looks constant."""
        board = [item(1, 10, date(2026, 8, 12)), item(2, 30)]
        curve = burndown_from_closures(
            board, "Sprint 1", today=date(2026, 8, 12)
        )
        # The 30-point item may have arrived on day 5; the curve cannot know.
        assert curve[0].remaining == 40.0


class TestStatusReconciliation:
    """Issue state is the fact; the board's Status column is a label."""

    def _args(self, **overrides):
        """Build a namespace for the reconciler.

        Args:
            **overrides: Fields to override.

        Returns:
            An argparse-like namespace.
        """
        import argparse

        base = {"issues_repo": None}
        base.update(overrides)
        return argparse.Namespace(**base)

    def test_closed_issue_counts_even_if_not_marked_done(self, monkeypatch, capsys):
        """An issue closed by a merged PR is complete regardless of Status."""
        from sprint_report import cli
        from sprint_report.models import ProjectItem

        monkeypatch.setattr(
            cli,
            "fetch_issues",
            lambda repo, since, **k: [
                {"number": 48, "state": "CLOSED", "closedAt": "2026-08-14T10:00:00Z"}
            ],
        )
        board = [
            ProjectItem(
                item_id="48",
                title="t",
                status="In Progress",
                iteration="Sprint 1",
                points=8,
                repository="acme/r",
            )
        ]
        assert not board[0].is_complete
        result = cli._apply_closure_dates(self._args(), board)
        assert result[0].is_complete
        assert result[0].closed_at == date(2026, 8, 14)
        assert "closed in GitHub but not marked Done" in capsys.readouterr().err

    def test_open_issue_stays_open(self, monkeypatch):
        """An issue still open is not promoted to complete."""
        from sprint_report import cli
        from sprint_report.models import ProjectItem

        monkeypatch.setattr(
            cli,
            "fetch_issues",
            lambda repo, since, **k: [{"number": 48, "state": "OPEN"}],
        )
        board = [
            ProjectItem(
                item_id="48", title="t", status="In Progress",
                iteration="Sprint 1", points=8, repository="acme/r",
            )
        ]
        assert not cli._apply_closure_dates(self._args(), board)[0].is_complete

    def test_done_on_board_without_a_github_issue(self, monkeypatch):
        """A draft item marked Done stays complete."""
        from sprint_report import cli
        from sprint_report.models import ProjectItem

        monkeypatch.setattr(cli, "fetch_issues", lambda repo, since, **k: [])
        board = [
            ProjectItem(
                item_id="1", title="draft", status="Done",
                iteration="Sprint 1", points=3, repository="acme/r",
            )
        ]
        assert cli._apply_closure_dates(self._args(), board)[0].is_complete

    def test_no_note_when_board_already_agrees(self, monkeypatch, capsys):
        """Nothing to reconcile means nothing to report."""
        from sprint_report import cli
        from sprint_report.models import ProjectItem

        monkeypatch.setattr(
            cli,
            "fetch_issues",
            lambda repo, since, **k: [
                {"number": 48, "state": "CLOSED", "closedAt": "2026-08-14T10:00:00Z"}
            ],
        )
        board = [
            ProjectItem(
                item_id="48", title="t", status="Done", iteration="Sprint 1",
                points=8, closed=True, repository="acme/r",
            )
        ]
        cli._apply_closure_dates(self._args(), board)
        assert "not marked Done" not in capsys.readouterr().err

    def test_unreadable_repo_degrades(self, monkeypatch, capsys):
        """A failed lookup warns and leaves the board data alone."""
        from sprint_report import cli
        from sprint_report.gh_source import GhError
        from sprint_report.models import ProjectItem

        def _boom(repo, since, **k):
            raise GhError("no access")

        monkeypatch.setattr(cli, "fetch_issues", _boom)
        board = [
            ProjectItem(
                item_id="1", title="t", status="Done", iteration="S1",
                points=3, repository="acme/r",
            )
        ]
        result = cli._apply_closure_dates(self._args(), board)
        assert result[0].is_complete
        assert "Issue state unavailable" in capsys.readouterr().err


class TestBurnup:
    """The burnup reads the same daily data the other way."""

    def test_completed_is_scope_minus_remaining(self):
        """The two framings are derived from one record, so cannot disagree."""
        from sprint_report.models import BurndownPoint

        point = BurndownPoint(date(2026, 8, 12), remaining=32.0, ideal=36.0, scope=40.0)
        assert point.completed == 8.0

    def test_completed_never_negative(self):
        """Scope shrinking below completed work is floored, not negative."""
        from sprint_report.models import BurndownPoint

        point = BurndownPoint(date(2026, 8, 12), remaining=50.0, ideal=0.0, scope=40.0)
        assert point.completed == 0.0

    def test_reconstructed_scope_is_flat(self):
        """Closure dates cannot say when an item joined, so scope is constant."""
        board = [
            item(1, 8, date(2026, 8, 12)),
            item(2, 5, date(2026, 8, 14)),
            item(3, 27),
        ]
        curve = burndown_from_closures(board, "Sprint 1", today=date(2026, 8, 16))
        assert {p.scope for p in curve} == {40.0}

    def test_reconstructed_completed_rises(self):
        """The burnup line climbs as work closes."""
        board = [item(1, 8, date(2026, 8, 12)), item(2, 32)]
        curve = burndown_from_closures(board, "Sprint 1", today=date(2026, 8, 14))
        assert curve[0].completed == 0.0
        assert curve[-1].completed == 8.0

    def test_snapshot_scope_tracks_changes(self):
        """Snapshots record scope per day, so growth becomes visible."""
        from sprint_report.metrics import burndown
        from sprint_report.models import Snapshot

        day_one = Snapshot(date(2026, 8, 10), "B", [item(1, 20)])
        day_two = Snapshot(date(2026, 8, 11), "B", [item(1, 20), item(2, 15)])
        curve = burndown([day_one, day_two], "Sprint 1")
        assert curve[0].scope == 20.0
        assert curve[1].scope == 35.0

    def test_deck_renders_a_burnup_slide(self, tmp_path):
        """A burndown always comes with its burnup."""
        from pptx import Presentation

        from sprint_report.deck import DeckBuilder
        from sprint_report.metrics import iteration_metrics

        board = [item(1, 8, date(2026, 8, 12)), item(2, 32)]
        curve = burndown_from_closures(board, "Sprint 1", today=date(2026, 8, 16))
        path = DeckBuilder("X").build(
            current=iteration_metrics(board, "Sprint 1"),
            history=[],
            burndown_points=curve,
            carryover=[],
            output_path=tmp_path / "burnup.pptx",
            burndown_reconstructed=True,
        )
        text = "\n".join(
            shape.text_frame.text
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Burnup" in text
        assert "Burndown" in text
        assert "scope is flat here" in text
