"""Tests for models, persistence, the API client, deck output, and the CLI.

The client tests stub the HTTP layer so no network access is required.
"""

from __future__ import annotations

import json
from datetime import date
from pathlib import Path

import pytest

from sprint_report.cli import _pick_iteration, _resolve_token, build_parser
from sprint_report.client import (
    FieldMapping,
    GitHubApiError,
    ProjectsClient,
)
from sprint_report.deck import DeckBuilder
from sprint_report.metrics import burndown, carryover_items, iteration_metrics
from sprint_report.models import ProjectItem, Snapshot
from sprint_report.snapshots import SnapshotStore


class TestProjectItem:
    """Behaviour of the item model."""

    def test_closed_issue_is_complete(self):
        """A closed issue counts as done regardless of board status."""
        assert ProjectItem(item_id="a", title="t", closed=True).is_complete

    def test_done_status_is_complete(self):
        """A Done status counts even when the issue is still open."""
        assert ProjectItem(item_id="a", title="t", status="Done").is_complete

    def test_status_matching_is_case_insensitive(self):
        """Status comparison tolerates casing differences."""
        assert ProjectItem(item_id="a", title="t", status="  DONE ").is_complete

    def test_in_progress_is_not_complete(self):
        """Work in flight is not counted as delivered."""
        assert not ProjectItem(item_id="a", title="t", status="In Progress").is_complete

    def test_missing_points_score_zero(self):
        """An unset estimate contributes nothing to a sum."""
        assert ProjectItem(item_id="a", title="t").effective_points == 0.0

    def test_iteration_end_is_inclusive(self):
        """A 14-day sprint starting the 10th ends on the 23rd."""
        item = ProjectItem(
            item_id="a",
            title="t",
            iteration="S1",
            iteration_start=date(2026, 8, 10),
            iteration_duration=14,
        )
        assert item.iteration_end == date(2026, 8, 23)

    def test_iteration_end_none_without_duration(self):
        """Without a duration there is no computable end date."""
        assert ProjectItem(item_id="a", title="t").iteration_end is None

    def test_json_round_trip(self):
        """Serialisation preserves every field including dates."""
        original = ProjectItem(
            item_id="a",
            title="t",
            iteration="S1",
            iteration_start=date(2026, 8, 10),
            iteration_duration=14,
            points=5.0,
            origin="Planned",
            closed_at=date(2026, 8, 15),
        )
        assert ProjectItem.from_json(original.to_json()) == original

    def test_from_json_ignores_unknown_keys(self):
        """Forward-compatibility: unrecognised keys are dropped, not fatal."""
        item = ProjectItem.from_json({"item_id": "a", "title": "t", "future_field": 1})
        assert item.item_id == "a"


class TestSnapshotStore:
    """Filesystem persistence of snapshots."""

    def test_write_then_load(self, tmp_path: Path, board_items):
        """A written snapshot reloads with the same item count."""
        store = SnapshotStore(tmp_path)
        store.write(Snapshot(date(2026, 8, 21), "Board", board_items))
        loaded = store.load_all()
        assert len(loaded) == 1
        assert loaded[0].item_count == len(board_items)

    def test_same_day_write_overwrites(self, tmp_path: Path, board_items):
        """Running twice in a day leaves one file, not two."""
        store = SnapshotStore(tmp_path)
        store.write(Snapshot(date(2026, 8, 21), "Board", board_items))
        store.write(Snapshot(date(2026, 8, 21), "Board", board_items[:2]))
        loaded = store.load_all()
        assert len(loaded) == 1
        assert loaded[0].item_count == 2

    def test_load_is_chronological(self, tmp_path: Path):
        """Snapshots return oldest first regardless of write order."""
        store = SnapshotStore(tmp_path)
        store.write(Snapshot(date(2026, 8, 21), "B", []))
        store.write(Snapshot(date(2026, 8, 19), "B", []))
        assert [snap.captured_on for snap in store.load_all()] == [
            date(2026, 8, 19),
            date(2026, 8, 21),
        ]

    def test_corrupt_file_is_skipped(self, tmp_path: Path):
        """One unreadable day does not break the whole run."""
        store = SnapshotStore(tmp_path)
        store.write(Snapshot(date(2026, 8, 21), "B", []))
        (tmp_path / "2026-08-20.json").write_text("{not json", encoding="utf-8")
        assert len(store.load_all()) == 1

    def test_missing_directory_is_created(self, tmp_path: Path):
        """The store creates its directory on demand."""
        target = tmp_path / "nested" / "snaps"
        SnapshotStore(target)
        assert target.exists()

    def test_load_for_iteration_filters(self, tmp_path: Path, snapshot_series):
        """Only snapshots mentioning the iteration are returned."""
        store = SnapshotStore(tmp_path)
        for snapshot in snapshot_series:
            store.write(snapshot)
        assert len(store.load_for_iteration("Sprint 14")) == 4
        assert store.load_for_iteration("Sprint 99") == []


class _StubResponse:
    """Minimal stand-in for a :class:`requests.Response`."""

    def __init__(self, payload: dict, status_code: int = 200) -> None:
        """Store the canned payload and status.

        Args:
            payload: Body returned by :meth:`json`.
            status_code: HTTP status to report.
        """
        self._payload = payload
        self.status_code = status_code
        self.text = json.dumps(payload)

    def json(self) -> dict:
        """Return the canned payload.

        Returns:
            The payload supplied at construction.
        """
        return self._payload


class TestProjectsClient:
    """GraphQL parsing, with the transport stubbed."""

    def test_requires_a_token(self):
        """Constructing without a token fails fast."""
        with pytest.raises(ValueError):
            ProjectsClient(token="")

    def test_graphql_errors_raise(self, monkeypatch):
        """API-level errors surface as :class:`GitHubApiError`."""
        client = ProjectsClient(token="x")
        monkeypatch.setattr(
            client._session,
            "post",
            lambda *a, **k: _StubResponse({"errors": [{"message": "nope"}]}),
        )
        with pytest.raises(GitHubApiError, match="nope"):
            client.resolve_project_id("acme", 1)

    def test_http_error_raises(self, monkeypatch):
        """A 401 is reported rather than retried indefinitely."""
        client = ProjectsClient(token="x")
        monkeypatch.setattr(
            client._session, "post", lambda *a, **k: _StubResponse({}, 401)
        )
        with pytest.raises(GitHubApiError, match="401"):
            client.resolve_project_id("acme", 1)

    def test_missing_project_raises(self, monkeypatch):
        """A valid response with no project is an error, not an empty board."""
        client = ProjectsClient(token="x")
        monkeypatch.setattr(
            client._session,
            "post",
            lambda *a, **k: _StubResponse(
                {"data": {"organization": {"projectV2": None}}}
            ),
        )
        with pytest.raises(GitHubApiError, match="No project"):
            client.resolve_project_id("acme", 1)

    def test_resolve_project_id(self, monkeypatch):
        """A well-formed response yields the node ID."""
        client = ProjectsClient(token="x")
        monkeypatch.setattr(
            client._session,
            "post",
            lambda *a, **k: _StubResponse(
                {"data": {"organization": {"projectV2": {"id": "PVT_1"}}}}
            ),
        )
        assert client.resolve_project_id("acme", 1) == "PVT_1"

    def test_parses_field_values(self):
        """Number, single-select, and iteration fields map onto the model."""
        client = ProjectsClient(token="x", fields=FieldMapping())
        item = client._parse_item(
            {
                "id": "PVTI_1",
                "content": {
                    "title": "Wire audit export",
                    "url": "https://example.test/1",
                    "state": "CLOSED",
                    "closedAt": "2026-08-15T10:00:00Z",
                    "repository": {"nameWithOwner": "acme/nxt"},
                    "milestone": {"title": "v2.0"},
                },
                "fieldValues": {
                    "nodes": [
                        {
                            "__typename": "ProjectV2ItemFieldNumberValue",
                            "number": 5,
                            "field": {"name": "Pts"},
                        },
                        {
                            "__typename": "ProjectV2ItemFieldSingleSelectValue",
                            "name": "Unplanned",
                            "field": {"name": "Origin"},
                        },
                        {
                            "__typename": "ProjectV2ItemFieldIterationValue",
                            "title": "Sprint 14",
                            "startDate": "2026-08-10",
                            "duration": 14,
                            "field": {"name": "Sprint"},
                        },
                    ]
                },
            }
        )
        assert item is not None
        assert item.points == 5
        assert item.origin == "Unplanned"
        assert item.iteration == "Sprint 14"
        assert item.iteration_start == date(2026, 8, 10)
        assert item.closed is True
        assert item.milestone == "v2.0"

    def test_custom_field_names_are_honoured(self):
        """Teams using different field names configure the mapping."""
        client = ProjectsClient(token="x", fields=FieldMapping(points="Story Points"))
        item = client._parse_item(
            {
                "id": "1",
                "content": {"title": "t"},
                "fieldValues": {
                    "nodes": [
                        {
                            "__typename": "ProjectV2ItemFieldNumberValue",
                            "number": 8,
                            "field": {"name": "Story Points"},
                        }
                    ]
                },
            }
        )
        assert item is not None and item.points == 8

    def test_item_without_content_is_skipped(self):
        """Items the token cannot see are dropped rather than half-parsed."""
        assert ProjectsClient(token="x")._parse_item({"id": "1", "content": None}) is None


class TestCli:
    """Argument handling and iteration selection."""

    def test_org_and_project_required_without_an_export(self, tmp_path: Path):
        """Identifying the board is mandatory unless an export is supplied."""
        from sprint_report.cli import _load_board

        args = build_parser().parse_args(["snapshot"])
        with pytest.raises(SystemExit, match="--org and --project"):
            _load_board(args)

    def test_export_needs_no_org_or_project(self, tmp_path: Path):
        """A saved export identifies the board on its own."""
        import json as _json

        from sprint_report.cli import _load_board

        export = tmp_path / "board.json"
        export.write_text(
            _json.dumps({"items": [{"id": "1", "title": "t"}]}), encoding="utf-8"
        )
        args = build_parser().parse_args(
            ["--from-export", str(export), "snapshot"]
        )
        title, items = _load_board(args)
        assert len(items) == 1

    def test_token_from_environment(self, monkeypatch):
        """The token falls back to the environment variable."""
        monkeypatch.setenv("GITHUB_PROJECTS_TOKEN", "env-token")
        assert _resolve_token(None) == "env-token"

    def test_missing_token_exits(self, monkeypatch):
        """No token anywhere is a hard stop with a useful message."""
        monkeypatch.delenv("GITHUB_PROJECTS_TOKEN", raising=False)
        with pytest.raises(SystemExit, match="No token"):
            _resolve_token(None)

    def test_explicit_iteration_selected(self, board_items):
        """A named iteration is used verbatim."""
        assert _pick_iteration(board_items, "Sprint 13") == "Sprint 13"

    def test_unknown_iteration_exits(self, board_items):
        """An unrecognised name lists what is available."""
        with pytest.raises(SystemExit, match="not found"):
            _pick_iteration(board_items, "Sprint 99")

    def test_current_falls_back_to_latest(self, board_items):
        """When no sprint spans today, the newest one is used."""
        assert _pick_iteration(board_items, "current") == "Sprint 14"

    def test_no_iterations_exits(self):
        """A board with no iterations cannot be reported on."""
        with pytest.raises(SystemExit, match="No iterations"):
            _pick_iteration([], "current")


class TestDeckBuilder:
    """End-to-end deck generation."""

    def _build(self, tmp_path: Path, board_items, snapshots=()):
        """Generate a deck from the fixtures.

        Args:
            tmp_path: pytest temporary directory.
            board_items: The synthetic board.
            snapshots: Optional snapshot series for the burndown.

        Returns:
            The path to the generated deck.
        """
        current = iteration_metrics(board_items, "Sprint 14")
        history = [
            iteration_metrics(board_items, title) for title in ("Sprint 12", "Sprint 13")
        ]
        return DeckBuilder("Acme Platform", "Acme Corp").build(
            current=current,
            history=history,
            burndown_points=burndown(list(snapshots), "Sprint 14"),
            carryover=carryover_items(board_items, "Sprint 14"),
            output_path=tmp_path / "review.pptx",
            milestone_forecasts=[("v2.0", 5.0)],
        )

    def test_writes_a_file(self, tmp_path: Path, board_items, snapshot_series):
        """The deck is created on disk."""
        assert self._build(tmp_path, board_items, snapshot_series).exists()

    def test_has_expected_slide_count(self, tmp_path: Path, board_items, snapshot_series):
        """Nine slides when history and snapshots are both present."""
        from pptx import Presentation

        path = self._build(tmp_path, board_items, snapshot_series)
        assert len(Presentation(str(path)).slides) == 9

    def test_survives_missing_snapshots(self, tmp_path: Path, board_items):
        """With no history the burndown slide degrades to a notice."""
        assert self._build(tmp_path, board_items).exists()

    def test_survives_empty_board(self, tmp_path: Path):
        """A board with nothing on it still produces a valid deck."""
        from sprint_report.models import SprintMetrics

        output = DeckBuilder("Empty").build(
            current=SprintMetrics("Sprint 1"),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "empty.pptx",
        )
        assert output.exists()

    def test_creates_parent_directory(self, tmp_path: Path, board_items):
        """A nested output path is created rather than failing."""
        from sprint_report.models import SprintMetrics

        output = DeckBuilder("X").build(
            current=SprintMetrics("Sprint 1"),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "deep" / "nested" / "deck.pptx",
        )
        assert output.exists()

    def test_carryover_table_truncates(self, tmp_path: Path):
        """More than eight carryover items are summarised, not overflowed."""
        from pptx import Presentation

        from sprint_report.models import SprintMetrics
        from tests.conftest import make_item

        many = [make_item(str(i), "S1", 3) for i in range(12)]
        output = DeckBuilder("X").build(
            current=SprintMetrics("S1"),
            history=[],
            burndown_points=[],
            carryover=many,
            output_path=tmp_path / "many.pptx",
        )
        text = "\n".join(
            shape.text_frame.text
            for slide in Presentation(str(output)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "further item(s) not shown" in text


class TestCliCommands:
    """End-to-end command execution with the API client stubbed."""

    @pytest.fixture()
    def stub_client(self, monkeypatch, board_items):
        """Patch the CLI's client factory to return a canned board.

        Args:
            monkeypatch: pytest monkeypatch fixture.
            board_items: The synthetic board fixture.

        Returns:
            The stub client instance used by the CLI.
        """
        from sprint_report import cli

        class _Stub:
            """Stand-in for :class:`ProjectsClient`."""

            def resolve_project_id(self, org: str, number: int) -> str:
                """Return a fixed node ID.

                Args:
                    org: Ignored.
                    number: Ignored.

                Returns:
                    A placeholder project node ID.
                """
                return "PVT_stub"

            def fetch_items(self, project_id: str):
                """Return the fixture board.

                Args:
                    project_id: Ignored.

                Returns:
                    A title and the fixture items.
                """
                return "Test Board", list(board_items)

        stub = _Stub()
        monkeypatch.setattr(
            cli, "_load_board", lambda args: ("Test Board", list(board_items))
        )
        return stub

    def test_snapshot_writes_file(self, stub_client, tmp_path, capsys):
        """The snapshot command persists today's board state."""
        from sprint_report.cli import main

        code = main(
            [
                "--org",
                "acme",
                "--project",
                "1",
                "--snapshots",
                str(tmp_path / "snaps"),
                "snapshot",
            ]
        )
        assert code == 0
        assert list((tmp_path / "snaps").glob("*.json"))
        assert "Captured 10 item(s)" in capsys.readouterr().out

    def test_report_writes_deck(self, stub_client, tmp_path, capsys):
        """The report command produces a deck for the chosen sprint."""
        from sprint_report.cli import main

        output = tmp_path / "review.pptx"
        code = main(
            [
                "--org",
                "acme",
                "--project",
                "1",
                "--snapshots",
                str(tmp_path / "snaps"),
                "report",
                "--iteration",
                "Sprint 14",
                "--output",
                str(output),
            ]
        )
        assert code == 0
        assert output.exists()

    def test_report_warns_without_snapshots(self, stub_client, tmp_path, capsys):
        """Missing burndown history is called out on stderr."""
        from sprint_report.cli import main

        main(
            [
                "--org",
                "acme",
                "--project",
                "1",
                "--snapshots",
                str(tmp_path / "snaps"),
                "report",
                "--iteration",
                "Sprint 14",
                "--output",
                str(tmp_path / "r.pptx"),
            ]
        )
        assert "no snapshots cover this sprint" in capsys.readouterr().err

    def test_api_error_returns_exit_code_one(self, monkeypatch, tmp_path):
        """An API failure is reported cleanly rather than traced."""
        from sprint_report import cli

        def _boom(args):
            raise cli.GitHubApiError("bad credentials")

        monkeypatch.setattr(cli, "_load_board", _boom)
        assert cli.main(["--org", "a", "--project", "1", "snapshot"]) == 1


class TestClientPagination:
    """Cursor handling and transient-failure retries."""

    def test_follows_pagination(self, monkeypatch):
        """All pages are fetched and concatenated."""
        client = ProjectsClient(token="x")
        pages = [
            {
                "data": {
                    "node": {
                        "title": "Board",
                        "items": {
                            "pageInfo": {"hasNextPage": True, "endCursor": "c1"},
                            "nodes": [
                                {
                                    "id": "1",
                                    "content": {"title": "a"},
                                    "fieldValues": {"nodes": []},
                                }
                            ],
                        },
                    }
                }
            },
            {
                "data": {
                    "node": {
                        "title": "Board",
                        "items": {
                            "pageInfo": {"hasNextPage": False, "endCursor": None},
                            "nodes": [
                                {
                                    "id": "2",
                                    "content": {"title": "b"},
                                    "fieldValues": {"nodes": []},
                                }
                            ],
                        },
                    }
                }
            },
        ]
        calls = iter(pages)
        monkeypatch.setattr(
            client._session, "post", lambda *a, **k: _StubResponse(next(calls))
        )
        title, items = client.fetch_items("PVT_1")
        assert title == "Board"
        assert [item.item_id for item in items] == ["1", "2"]

    def test_missing_node_raises(self, monkeypatch):
        """A board the token cannot see is an error, not an empty result."""
        client = ProjectsClient(token="x")
        monkeypatch.setattr(
            client._session,
            "post",
            lambda *a, **k: _StubResponse({"data": {"node": None}}),
        )
        with pytest.raises(GitHubApiError, match="not found or not visible"):
            client.fetch_items("PVT_1")

    def test_retries_transient_failure(self, monkeypatch):
        """A 503 is retried before succeeding."""
        client = ProjectsClient(token="x", max_retries=3)
        monkeypatch.setattr("time.sleep", lambda _: None)
        responses = iter(
            [
                _StubResponse({}, 503),
                _StubResponse({"data": {"organization": {"projectV2": {"id": "PVT_9"}}}}),
            ]
        )
        monkeypatch.setattr(client._session, "post", lambda *a, **k: next(responses))
        assert client.resolve_project_id("acme", 1) == "PVT_9"

    def test_gives_up_after_max_retries(self, monkeypatch):
        """Persistent failure surfaces rather than looping forever."""
        client = ProjectsClient(token="x", max_retries=2)
        monkeypatch.setattr("time.sleep", lambda _: None)
        monkeypatch.setattr(
            client._session, "post", lambda *a, **k: _StubResponse({}, 502)
        )
        with pytest.raises(GitHubApiError, match="Giving up"):
            client.resolve_project_id("acme", 1)


class TestLinkedTableSlides:
    """The unestimated and unsprinted follow-up lists."""

    def _build(self, tmp_path: Path, **kwargs):
        """Generate a deck with the follow-up lists populated.

        Args:
            tmp_path: pytest temporary directory.
            **kwargs: Passed through to :meth:`DeckBuilder.build`.

        Returns:
            The generated deck path.
        """
        from sprint_report.models import SprintMetrics

        return DeckBuilder("X").build(
            current=SprintMetrics("Sprint 1"),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "deck.pptx",
            **kwargs,
        )

    def _text(self, path: Path) -> str:
        """Flatten all slide text, including table cells.

        Args:
            path: Deck to read.

        Returns:
            Newline-joined text of every shape and table cell.
        """
        from pptx import Presentation

        chunks = []
        for slide in Presentation(str(path)).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            chunks.append(cell.text)
        return "\n".join(chunks)

    def test_lists_unestimated_items(self, tmp_path: Path):
        """Unestimated items appear by number and title."""
        item = ProjectItem(
            item_id="48", title="PHI-access audit log", status="Done"
        )
        text = self._text(self._build(tmp_path, unestimated=[item]))
        assert "#48" in text
        assert "PHI-access audit log" in text

    def test_lists_unsprinted_items(self, tmp_path: Path):
        """Items outside every iteration appear on their own slide."""
        item = ProjectItem(item_id="21", title="Routes", status="Todo")
        assert "Routes" in self._text(self._build(tmp_path, unsprinted=[item]))

    def test_titles_carry_hyperlinks(self, tmp_path: Path):
        """A row's title links back to the issue."""
        from pptx import Presentation

        item = ProjectItem(
            item_id="21",
            title="Routes",
            url="https://github.com/acme/repo/issues/21",
            status="Todo",
        )
        path = self._build(tmp_path, unsprinted=[item])
        addresses = [
            run.hyperlink.address
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_table
            for row in shape.table.rows
            for cell in row.cells
            for paragraph in cell.text_frame.paragraphs
            for run in paragraph.runs
            if run.hyperlink.address
        ]
        assert "https://github.com/acme/repo/issues/21" in addresses

    def test_empty_lists_render_a_notice(self, tmp_path: Path):
        """Nothing to follow up on is stated rather than left blank."""
        text = self._text(self._build(tmp_path))
        assert "Every item in this sprint carries an estimate." in text
        assert "Every board item is assigned to an iteration." in text

    def test_long_lists_are_truncated(self, tmp_path: Path):
        """Beyond the row cap the remainder is summarised."""
        many = [
            ProjectItem(item_id=str(i), title=f"Item {i}") for i in range(20)
        ]
        text = self._text(self._build(tmp_path, unsprinted=many))
        assert "further item(s) not shown" in text

    def test_long_titles_are_elided(self, tmp_path: Path):
        """Overlong titles are shortened rather than overflowing the cell."""
        item = ProjectItem(item_id="1", title="x" * 200)
        assert "…" in self._text(self._build(tmp_path, unestimated=[item]))


class TestTrendSlideGuards:
    """Trend slides are omitted rather than rendered with one data point."""

    def _slide_count(self, tmp_path: Path, **kwargs) -> int:
        """Build a deck and count its slides.

        Args:
            tmp_path: pytest temporary directory.
            **kwargs: Passed through to :meth:`DeckBuilder.build`.

        Returns:
            Number of slides in the generated deck.
        """
        from pptx import Presentation

        from sprint_report.models import SprintMetrics

        path = DeckBuilder("X").build(
            current=SprintMetrics("Sprint 1", committed_points=34),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "deck.pptx",
            **kwargs,
        )
        return len(Presentation(str(path)).slides)

    def test_first_sprint_omits_trend_slides(self, tmp_path: Path):
        """No prior sprints means no trend, but work mix still renders."""
        assert self._slide_count(tmp_path) == 6

    def test_work_mix_renders_without_history(self, tmp_path: Path):
        """A single sprint's composition is meaningful on its own."""
        from pptx import Presentation

        from sprint_report.models import SprintMetrics

        path = DeckBuilder("X").build(
            current=SprintMetrics(
                "Sprint 1",
                committed_points=110,
                planned_points=83,
                unplanned_points=22,
                carryover_points=5,
            ),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "mix.pptx",
        )
        text = "\n".join(
            shape.text_frame.text
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Work mix" in text
        assert "What this sprint was actually made of" in text

    def test_work_mix_omitted_when_nothing_committed(self, tmp_path: Path):
        """An empty sprint has no composition to show."""
        from pptx import Presentation

        from sprint_report.models import SprintMetrics

        path = DeckBuilder("X").build(
            current=SprintMetrics("Sprint 1"),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "empty.pptx",
        )
        text = "\n".join(
            shape.text_frame.text
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Work mix" not in text

    def test_no_milestones_omits_forecast(self, tmp_path: Path):
        """An empty forecast is left out rather than shown as a notice."""
        from pptx import Presentation

        from sprint_report.models import SprintMetrics

        path = DeckBuilder("X").build(
            current=SprintMetrics("S1"),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "nf.pptx",
        )
        text = "\n".join(
            shape.text_frame.text
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Forecast" not in text

    def test_first_sprint_omits_burndown(self, tmp_path: Path):
        """No snapshots means no burndown slide at all, not an empty one."""
        from pptx import Presentation

        from sprint_report.models import SprintMetrics

        path = DeckBuilder("X").build(
            current=SprintMetrics("Sprint 1"),
            history=[],
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "d.pptx",
        )
        text = "\n".join(
            shape.text_frame.text
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Burndown" not in text

    def test_history_restores_trend_slides(self, tmp_path: Path, board_items):
        """Once a prior sprint exists the trend slides come back."""
        from sprint_report.metrics import iteration_metrics

        history = [iteration_metrics(board_items, "Sprint 12")]
        assert self._slide_count(tmp_path) == 6
        from pptx import Presentation

        from sprint_report.models import SprintMetrics

        path = DeckBuilder("X").build(
            current=SprintMetrics("Sprint 13", committed_points=20),
            history=history,
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "with-history.pptx",
        )
        assert len(Presentation(str(path)).slides) == 7


class TestModeAndSummary:
    """Mid-sprint mode selection and the machine-readable summary."""

    def test_explicit_mode_wins(self):
        """An explicit --mode is never overridden by the date."""
        from datetime import date

        from sprint_report.cli import _resolve_mode
        from sprint_report.models import SprintMetrics

        metrics = SprintMetrics("S1", date(2026, 8, 10), date(2026, 8, 23))
        assert _resolve_mode("review", metrics) == "review"
        assert _resolve_mode("midsprint", metrics) == "midsprint"

    def test_auto_picks_review_for_a_closed_sprint(self):
        """A sprint that ended in the past gets a review deck."""
        from datetime import date

        from sprint_report.cli import _resolve_mode
        from sprint_report.models import SprintMetrics

        metrics = SprintMetrics("S1", date(2020, 1, 1), date(2020, 1, 14))
        assert _resolve_mode("auto", metrics) == "review"

    def test_auto_picks_midsprint_inside_the_window(self, monkeypatch):
        """A sprint spanning today gets a mid-sprint deck."""
        from datetime import date

        from sprint_report import cli
        from sprint_report.models import SprintMetrics

        monkeypatch.setattr(cli, "utc_today", lambda: date(2026, 8, 17))
        metrics = SprintMetrics("S1", date(2026, 8, 10), date(2026, 8, 23))
        assert cli._resolve_mode("auto", metrics) == "midsprint"

    def test_auto_without_dates_falls_back_to_review(self):
        """An iteration with no dates cannot be mid-flight."""
        from sprint_report.cli import _resolve_mode
        from sprint_report.models import SprintMetrics

        assert _resolve_mode("auto", SprintMetrics("S1")) == "review"

    def test_unknown_deck_mode_rejected(self):
        """A typo in the mode is a programming error, not a silent default."""
        with pytest.raises(ValueError, match="Unknown mode"):
            DeckBuilder("X", mode="weekly")

    def test_midsprint_relabels_the_open_work_slide(self, tmp_path: Path):
        """Mid-sprint decks say 'In flight', not 'Carryover'."""
        from pptx import Presentation

        from sprint_report.models import SprintMetrics

        path = DeckBuilder("X", mode="midsprint").build(
            current=SprintMetrics("S1"),
            history=[],
            burndown_points=[],
            carryover=[ProjectItem(item_id="1", title="t", points=3)],
            output_path=tmp_path / "m.pptx",
        )
        text = "\n".join(
            shape.text_frame.text
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "In flight" in text
        assert "points still open" in text
        assert "rolling forward" not in text

    def test_summary_reports_estimate_coverage(self):
        """Coverage is the share of items that carry an estimate."""
        from sprint_report.cli import _summary_payload
        from sprint_report.models import SprintMetrics

        metrics = SprintMetrics("S1", total_items=16, completed_items=10)
        unestimated = [ProjectItem(item_id=str(i), title="t") for i in range(9)]
        payload = _summary_payload(metrics, "review", unestimated, [], [], False)
        assert payload["data_quality"]["estimate_coverage_percent"] == 43.8
        assert payload["data_quality"]["unestimated_items"] == 9
        assert payload["items"]["percent_complete"] == 62.5

    def test_summary_flags_unused_origin_field(self):
        """All-planned work means the Origin field is not being maintained."""
        from sprint_report.cli import _summary_payload
        from sprint_report.models import SprintMetrics

        payload = _summary_payload(
            SprintMetrics("S1", planned_points=34), "review", [], [], [], False
        )
        assert payload["data_quality"]["origin_field_used"] is False

    def test_summary_survives_an_empty_sprint(self):
        """No items means zeroes, not a division error."""
        from sprint_report.cli import _summary_payload
        from sprint_report.models import SprintMetrics

        payload = _summary_payload(SprintMetrics("S1"), "review", [], [], [], False)
        assert payload["items"]["percent_complete"] == 0.0
        assert payload["data_quality"]["estimate_coverage_percent"] == 0.0

    def test_summary_is_json_serialisable(self):
        """The payload must survive json.dumps unchanged."""
        import json

        from sprint_report.cli import _summary_payload
        from sprint_report.models import SprintMetrics

        payload = _summary_payload(
            SprintMetrics("S1", total_items=1),
            "review",
            [],
            [],
            [ProjectItem(item_id="1", title="t", url="https://x/1", points=3)],
            True,
        )
        assert json.loads(json.dumps(payload))["open_items"][0]["id"] == "1"


class TestTrendSlide:
    """Velocity and predictability are one chart, not two slides."""

    def _deck(self, tmp_path: Path, board_items):
        """Build a deck with two sprints of history.

        Args:
            tmp_path: pytest temporary directory.
            board_items: The synthetic board fixture.

        Returns:
            The generated deck path.
        """
        from sprint_report.metrics import iteration_metrics
        from sprint_report.models import SprintMetrics

        history = [
            iteration_metrics(board_items, "Sprint 12"),
            iteration_metrics(board_items, "Sprint 13"),
        ]
        return DeckBuilder("X").build(
            current=SprintMetrics("Sprint 14", committed_points=16,
                                  completed_points=11),
            history=history,
            burndown_points=[],
            carryover=[],
            output_path=tmp_path / "trend.pptx",
        )

    def _text(self, path: Path) -> str:
        """Flatten slide text.

        Args:
            path: Deck to read.

        Returns:
            All shape text joined by newlines.
        """
        from pptx import Presentation

        return "\n".join(
            shape.text_frame.text
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )

    def test_no_separate_velocity_slide(self, tmp_path: Path, board_items):
        """Velocity is the completed series, not its own slide."""
        text = self._text(self._deck(tmp_path, board_items))
        assert "Velocity" not in text
        assert "Delivery trend" in text

    def test_no_separate_predictability_slide(self, tmp_path: Path, board_items):
        """Predictability is folded into the same chart."""
        assert "Predictability" not in self._text(self._deck(tmp_path, board_items))

    def test_shows_ratios_per_sprint(self, tmp_path: Path, board_items):
        """The delivered-against-commitment figures appear as text."""
        text = self._text(self._deck(tmp_path, board_items))
        assert "Delivered against commitment" in text

    def test_columns_are_overlapped_not_clustered(self, tmp_path: Path, board_items):
        """Full overlap makes each sprint read as one progress bar."""
        from pptx import Presentation

        path = self._deck(tmp_path, board_items)
        plots = [
            shape.chart.plots[0]
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.has_chart
        ]
        assert plots[0].overlap == 100

    def test_work_mix_stays_stacked(self, tmp_path: Path, board_items):
        """The planned/unplanned/carryover breakdown remains a stack."""
        assert "Work mix" in self._text(self._deck(tmp_path, board_items))


class TestModeAtSprintBoundary:
    """The final day of a sprint is review day."""

    def test_last_day_is_a_review(self, monkeypatch):
        """A review held on the closing day should not say 'mid-sprint'."""
        from datetime import date

        from sprint_report import cli
        from sprint_report.models import SprintMetrics

        monkeypatch.setattr(cli, "utc_today", lambda: date(2026, 8, 23))
        metrics = SprintMetrics("Sprint 1", date(2026, 8, 10), date(2026, 8, 23))
        assert cli._resolve_mode("auto", metrics) == "review"

    def test_day_before_the_end_is_midsprint(self, monkeypatch):
        """The day before close is still work in flight."""
        from datetime import date

        from sprint_report import cli
        from sprint_report.models import SprintMetrics

        monkeypatch.setattr(cli, "utc_today", lambda: date(2026, 8, 22))
        metrics = SprintMetrics("Sprint 1", date(2026, 8, 10), date(2026, 8, 23))
        assert cli._resolve_mode("auto", metrics) == "midsprint"

    def test_first_day_is_midsprint(self, monkeypatch):
        """Day one is not a review."""
        from datetime import date

        from sprint_report import cli
        from sprint_report.models import SprintMetrics

        monkeypatch.setattr(cli, "utc_today", lambda: date(2026, 8, 10))
        metrics = SprintMetrics("Sprint 1", date(2026, 8, 10), date(2026, 8, 23))
        assert cli._resolve_mode("auto", metrics) == "midsprint"

    def test_after_the_end_is_a_review(self, monkeypatch):
        """A closed sprint always reviews."""
        from datetime import date

        from sprint_report import cli
        from sprint_report.models import SprintMetrics

        monkeypatch.setattr(cli, "utc_today", lambda: date(2026, 8, 30))
        metrics = SprintMetrics("Sprint 1", date(2026, 8, 10), date(2026, 8, 23))
        assert cli._resolve_mode("auto", metrics) == "review"
